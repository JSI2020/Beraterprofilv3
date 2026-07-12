"""Compare paragraph indent XML between template and output."""

import glob
import sys
from pptx import Presentation
from lxml import etree

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

def ppr_info(p):
    ppr = p._p.find("a:pPr", NS)
    if ppr is None:
        return {}
    info = {"level": p.level}
    for tag in ("marL", "indent", "lvl"):
        el = ppr.get(tag)
        if el:
            info[tag] = el
    bu = ppr.find("a:buChar", NS)
    if bu is not None:
        info["buChar"] = bu.get("char")
    buauto = ppr.find("a:buAutoNum", NS)
    if buauto is not None:
        info["buAutoNum"] = buauto.get("type")
    bunone = ppr.find("a:buNone", NS)
    if bunone is not None:
        info["buNone"] = True
    return info


def dump(path, label):
    print(f"\n######## {label}: {path}")
    prs = Presentation(path)
    slide = prs.slides[0]
    regions = {
        "kompetenzen": (533999, 2943147),
        "tools": (553881, 4846422),
        "abschluss": (8022387, 4720513),
        "title": (503148, 27489),
        "schwerpunkte": (3727239, 1291395),
    }
    for name, (tl, tt) in regions.items():
        best, bd = None, 10**18
        for s in slide.shapes:
            if not s.has_text_frame:
                continue
            d = abs(s.left - tl) + abs(s.top - tt)
            if d < bd and d <= 200000:
                bd, best = d, s
        if not best:
            continue
        print(f"\n--- {name} ---")
        for i, p in enumerate(best.text_frame.paragraphs):
            if p.text.strip():
                print(f"  p{i} L{p.level} {ppr_info(p)} | {p.text[:60]}")


dump("templates/beraterprofil_template.pptx", "TEMPLATE")
out = sys.argv[1] if len(sys.argv) > 1 else sorted(glob.glob("output/*.pptx"))[-1]
dump(out, "OUTPUT")

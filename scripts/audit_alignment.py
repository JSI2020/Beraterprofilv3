"""Full alignment audit: template vs output."""

import glob
import sys
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

EMU_PER_INCH = 914400


def align_name(p):
    if p.alignment is None:
        return "inherit"
    names = {PP_ALIGN.LEFT: "LEFT", PP_ALIGN.CENTER: "CENTER", PP_ALIGN.RIGHT: "RIGHT", PP_ALIGN.JUSTIFY: "JUSTIFY"}
    return names.get(p.alignment, str(p.alignment))


def audit_shape(shape, name):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    print(f"\n  [{name}] pos=({shape.left},{shape.top}) size=({shape.width}x{shape.height})")
    print(f"    tf margins LRTB=({tf.margin_left},{tf.margin_right},{tf.margin_top},{tf.margin_bottom})")
    print(f"    word_wrap={tf.word_wrap} auto_size={tf.auto_size}")
    for i, p in enumerate(tf.paragraphs):
        if not p.text.strip():
            continue
        print(f"    p{i} lvl={p.level} align={align_name(p)} spcB={p.space_before} spcA={p.space_after}")
        print(f"         text={p.text[:90]!r}")


def audit_file(path, label):
    print(f"\n{'='*70}\n{label}: {path}\n{'='*70}")
    prs = Presentation(path)
    slide = prs.slides[0]
    # Header area shapes by position
    targets = [
        ("title", 503148, 27489),
        ("photo", 503148, 874880),
        ("labels_pos_schwerp", 2209140, 972000),
        ("label_summary", 2209140, 1569297),
        ("position_val", 3727239, 970730),
        ("schwerpunkte_val", 3727239, 1291395),
        ("summary_val", 3727239, 1590943),
        ("kompetenzen", 533999, 2943147),
        ("erfahrungen", 4062770, 2906571),
        ("ausbildung", 7997235, 2906571),
        ("abschluss", 8022387, 4720513),
        ("tools", 553881, 4846422),
    ]
    for name, tl, tt in targets:
        best, bd = None, 10**18
        for s in slide.shapes:
            if not s.has_text_frame:
                continue
            d = abs(s.left - tl) + abs(s.top - tt)
            if d < bd and d <= 250000:
                bd, best = d, s
        if best:
            audit_shape(best, f"{name} (d={bd})")


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else sorted(glob.glob("output/*.pptx"))[-1]
    audit_file("templates/beraterprofil_template.pptx", "TEMPLATE")
    audit_file(out, "OUTPUT")

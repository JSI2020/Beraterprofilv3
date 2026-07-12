"""Compare template reference content lengths vs generated output."""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent


def analyze(path: Path) -> dict:
    prs = Presentation(str(path))
    slide = prs.slides[0]
    regions = {
        "title": (503148, 27489),
        "kompetenzen": (533999, 2943147),
        "relevante_erfahrungen": (4062770, 2906571),
        "ausbildung_karriere": (7997235, 2906571),
        "position": (3727239, 1050858),
        "schwerpunkte": (3727239, 1363311),
        "summary": (3727239, 1590943),
        "abschluss_zertifikate": (8022387, 4720513),
        "tool_kenntnisse": (553881, 4846422),
    }
    tol = 250000
    result = {}
    for key, (tl, tt) in regions.items():
        best = None
        best_d = 1e18
        for shape in slide.shapes:
            d = abs(shape.left - tl) + abs(shape.top - tt)
            if d < best_d and d <= tol and shape.has_text_frame:
                best_d = d
                best = shape
        if not best:
            continue
        tf = best.text_frame
        paras = []
        for p in tf.paragraphs:
            if p.text.strip():
                paras.append(
                    {
                        "text": p.text,
                        "len": len(p.text),
                        "level": p.level,
                        "word_wrap": tf.word_wrap,
                        "auto_size": str(tf.auto_size),
                        "margin_l": tf.margin_left,
                        "margin_r": tf.margin_right,
                        "margin_t": tf.margin_top,
                        "margin_b": tf.margin_bottom,
                    }
                )
        result[key] = {
            "width_emu": best.width,
            "height_emu": best.height,
            "width_in": round(best.width / 914400, 2),
            "height_in": round(best.height / 914400, 2),
            "paragraph_count": len([p for p in tf.paragraphs if p.text.strip()]),
            "total_chars": sum(len(p.text) for p in tf.paragraphs),
            "paragraphs": paras,
        }
    return result


if __name__ == "__main__":
    for p in sys.argv[1:]:
        data = analyze(Path(p))
        print(f"\n=== {p} ===")
        print(json.dumps(data, indent=2, ensure_ascii=False))

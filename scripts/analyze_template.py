"""One-off analysis script for Beraterprofil template structure."""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def shape_info(shape, depth=0):
    info = {
        "name": shape.name,
        "type": str(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }
    if shape.has_text_frame:
        paras = []
        for p in shape.text_frame.paragraphs:
            runs = [r.text for r in p.runs]
            paras.append(
                {
                    "text": p.text,
                    "level": p.level,
                    "alignment": str(p.alignment) if p.alignment else None,
                    "runs": runs,
                    "font_sizes": [
                        r.font.size.pt if r.font.size else None for r in p.runs
                    ],
                    "bold": [r.font.bold for r in p.runs],
                }
            )
        info["paragraphs"] = paras
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["children"] = [shape_info(s, depth + 1) for s in shape.shapes]
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["picture"] = True
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "cells": [
                [table.cell(r, c).text for c in range(len(table.columns))]
                for r in range(len(table.rows))
            ],
        }
    return info


def main(pptx_path: str) -> None:
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"Slide count: {len(prs.slides)}")
    shapes = []
    for i, shape in enumerate(slide.shapes):
        s = shape_info(shape)
        s["index"] = i
        shapes.append(s)
    out = Path(__file__).parent / "template_analysis.json"
    out.write_text(json.dumps(shapes, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Wrote {out}")
    for s in shapes:
        if s.get("paragraphs"):
            texts = [p["text"] for p in s["paragraphs"] if p["text"].strip()]
            if texts:
                print(f"\n--- Shape {s['index']} ({s['name']}) ---")
                for t in texts:
                    print(t[:200])


if __name__ == "__main__":
    main(sys.argv[1])

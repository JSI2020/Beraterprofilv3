"""Inspect bullet/paragraph XML in template shapes."""

from pathlib import Path
from pptx import Presentation
from lxml import etree

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

prs = Presentation("templates/beraterprofil_template.pptx")
slide = prs.slides[0]

regions = {
    "kompetenzen": (533999, 2943147),
    "erfahrungen": (4062770, 2906571),
    "ausbildung": (7997235, 2906571),
    "abschluss": (8022387, 4720513),
    "tools": (553881, 4846422),
}

for name, (tl, tt) in regions.items():
    for shape in slide.shapes:
        if abs(shape.left - tl) + abs(shape.top - tt) > 200000:
            continue
        print(f"\n=== {name} ===")
        tf = shape.text_frame
        for i, p in enumerate(tf.paragraphs):
            if not p.text.strip():
                continue
            lvl = p.level
            ppr = p._p.find("a:pPr", NS)
            bu = None
            if ppr is not None:
                bu = ppr.find("a:buChar", NS)
                bun = ppr.find("a:buNone", NS)
                buauto = ppr.find("a:buAutoNum", NS)
            bu_info = "none"
            if bu is not None:
                bu_info = f"char={bu.get('char')}"
            elif bun is not None:
                bu_info = "buNone"
            elif buauto is not None:
                bu_info = "autoNum"
            print(f"  p{i} L{lvl} bu={bu_info} | {p.text[:70]}")
        break

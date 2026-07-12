from pathlib import Path
import glob
from pptx import Presentation
from lxml import etree

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
f = sorted(glob.glob("output/*.pptx"))[-1]
# Allow: python verify_output.py output/specific.pptx
import sys
if len(sys.argv) > 1:
    f = sys.argv[1]
print("FILE:", f)
prs = Presentation(f)
slide = prs.slides[0]
regions = {
    "kompetenzen": (533999, 2943147),
    "erfahrungen": (4062770, 2906571),
    "ausbildung": (7997235, 2906571),
    "abschluss": (8022387, 4720513),
    "tools": (553881, 4846422),
}
for name, (tl, tt) in regions.items():
    best = None
    best_d = 10**18
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        d = abs(s.left - tl) + abs(s.top - tt)
        if d < best_d and d <= 200000:
            best_d = d
            best = s
    if best is None:
        continue
    print(f"\n=== {name} (d={best_d}) ===")
    for p in best.text_frame.paragraphs:
        if p.text.strip():
            ppr = p._p.find("a:pPr", NS)
            bu = ppr.find("a:buChar", NS) if ppr is not None else None
            buc = bu.get("char") if bu is not None else "none"
            print(f"  L{p.level} bu={buc} | {p.text[:80]}")

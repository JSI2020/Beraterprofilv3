"""Dump all shapes from a PPTX for reference comparison."""

from pathlib import Path
import sys
from pptx import Presentation

path = sys.argv[1] if len(sys.argv) > 1 else "templates/beraterprofil_reference_example.pptx"
prs = Presentation(path)
slide = prs.slides[0]
print(f"FILE: {path}")
print(f"Slide size: {prs.slide_width} x {prs.slide_height}")
for i, s in enumerate(slide.shapes):
    t = ""
    if s.has_text_frame:
        t = s.text_frame.text.replace("\n", " | ")[:100]
    print(f"{i:2d} {s.name:32s} L={s.left:8d} T={s.top:8d} W={s.width:8d} H={s.height:7d} | {t}")

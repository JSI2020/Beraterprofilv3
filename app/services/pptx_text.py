"""Low-level PowerPoint text helpers — preserve template paragraph formatting."""

from __future__ import annotations

from copy import deepcopy

from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement


def _tx_body(text_frame):
    return text_frame._txBody


def _run_properties(run_element) -> OxmlElement | None:
    r_pr = run_element.find(qn("a:rPr"))
    return deepcopy(r_pr) if r_pr is not None else None


def _make_run(text: str, r_pr: OxmlElement | None = None) -> OxmlElement:
    run = OxmlElement("a:r")
    if r_pr is not None:
        run.append(deepcopy(r_pr))
    text_element = OxmlElement("a:t")
    if text.startswith(" ") or text.endswith(" "):
        text_element.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    text_element.text = text
    run.append(text_element)
    return run


def _replace_runs(paragraph, runs: list[OxmlElement]) -> None:
    for run in list(paragraph.findall(qn("a:r"))):
        paragraph.remove(run)
    end_para_rpr = paragraph.find(qn("a:endParaRPr"))
    for run in runs:
        if end_para_rpr is not None:
            end_para_rpr.addprevious(run)
        else:
            paragraph.append(run)


def _trim_paragraphs(text_frame, keep: int) -> None:
    body = _tx_body(text_frame)
    paragraphs = body.findall(qn("a:p"))
    for paragraph in paragraphs[keep:]:
        body.remove(paragraph)


def set_minimal_plain_text(text_frame, text: str) -> None:
    if not text_frame.paragraphs:
        return
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text
    _trim_paragraphs(text_frame, 1)


def set_minimal_bullets(text_frame, lines: list[str]) -> None:
    if not lines:
        return
    paragraphs = text_frame.paragraphs
    for index, line in enumerate(lines):
        if index >= len(paragraphs):
            break
        paragraph = paragraphs[index]
        if paragraph.runs:
            paragraph.runs[0].text = line
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = line
    for index in range(len(lines), len(paragraphs)):
        paragraph = paragraphs[index]
        if paragraph.runs:
            paragraph.runs[0].text = ""
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = ""


def set_plain_text(text_frame, text: str, prototype) -> None:
    proto_runs = prototype.findall(qn("a:r"))
    r_pr = _run_properties(proto_runs[0]) if proto_runs else None
    paragraphs = text_frame.paragraphs
    if paragraphs:
        _replace_runs(paragraphs[0]._p, [_make_run(text, r_pr)])
        _trim_paragraphs(text_frame, 1)
    else:
        paragraph = deepcopy(prototype)
        _replace_runs(paragraph, [_make_run(text, r_pr)])
        _tx_body(text_frame).append(paragraph)


def set_bullet_lines(text_frame, lines: list[str], prototype) -> None:
    if not lines:
        return
    proto_runs = prototype.findall(qn("a:r"))
    r_pr = _run_properties(proto_runs[0]) if proto_runs else None
    existing = list(text_frame.paragraphs)
    for index, line in enumerate(lines):
        if index < len(existing):
            paragraph = existing[index]._p
        else:
            paragraph = deepcopy(prototype)
            _tx_body(text_frame).append(paragraph)
        _replace_runs(paragraph, [_make_run(line, r_pr)])
    _trim_paragraphs(text_frame, len(lines))


def set_categorized_lines(text_frame, items: list[tuple[str, str]], prototype) -> None:
    if not items:
        return
    proto_runs = prototype.findall(qn("a:r"))
    bold_r_pr = _run_properties(proto_runs[0]) if proto_runs else None
    normal_r_pr = _run_properties(proto_runs[-1]) if proto_runs else bold_r_pr
    existing = list(text_frame.paragraphs)
    for index, (category, details) in enumerate(items):
        if index < len(existing):
            paragraph = existing[index]._p
        else:
            paragraph = deepcopy(prototype)
            _tx_body(text_frame).append(paragraph)
        _replace_runs(
            paragraph,
            [
                _make_run(category, bold_r_pr),
                _make_run(f": {details}", normal_r_pr),
            ],
        )
    _trim_paragraphs(text_frame, len(items))


def normalize_paragraph_properties(text_frame) -> None:
    """Copy first paragraph pPr to later lines that lost marL/indent/buChar."""
    paragraphs = text_frame.paragraphs
    if len(paragraphs) < 2:
        return
    source_ppr = paragraphs[0]._p.find(qn("a:pPr"))
    if source_ppr is None:
        return
    for paragraph in paragraphs[1:]:
        target = paragraph._p
        old = target.find(qn("a:pPr"))
        if old is not None:
            target.remove(old)
        target.insert(0, deepcopy(source_ppr))

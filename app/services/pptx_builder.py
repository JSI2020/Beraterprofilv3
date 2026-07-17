"""Fill ORBIT Beraterprofil template — match example2 Textfeld layout."""

from __future__ import annotations

import io
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

from app.schemas.profile import BeraterprofilData
from app.services.content_fit import fit_profile
from app.services.pptx_shapes import clear_shape_text, find_all_content_shapes
from app.services.pptx_text import (
    set_bullet_lines,
    set_categorized_lines,
    set_minimal_plain_text,
    set_plain_text,
)


def _prototype(shape):
    return deepcopy(shape.text_frame.paragraphs[0]._p)


def _replace_photo(slide, photo_bytes: bytes, shapes: dict) -> None:
    shape = shapes["photo"]
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(io.BytesIO(photo_bytes), left, top, width=width, height=height)


def build_pptx(
    data: BeraterprofilData,
    template_path: Path,
    output_path: Path,
    photo_bytes: bytes | None = None,
) -> Path:
    # Profile in UI/JSON is already normalized — only apply length limits, do not re-split sections.
    data = fit_profile(data, preserve_sections=True)

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    shapes = find_all_content_shapes(slide)

    prototypes = {
        "title": _prototype(shapes["title"]),
        "position": _prototype(shapes["position"]),
        "kompetenzen": _prototype(shapes["kompetenzen"]),
        "relevante_erfahrungen": _prototype(shapes["relevante_erfahrungen"]),
        "ausbildung_karriere": _prototype(shapes["ausbildung_karriere"]),
        "abschluss_zertifikate": _prototype(shapes["abschluss_zertifikate"]),
        "tool_kenntnisse": _prototype(shapes["tool_kenntnisse"]),
    }

    set_plain_text(
        shapes["title"].text_frame,
        f"Beraterprofil – {data.title_domain}",
        prototypes["title"],
    )
    set_plain_text(shapes["position"].text_frame, data.position, prototypes["position"])
    set_minimal_plain_text(shapes["schwerpunkte"].text_frame, data.schwerpunkte)
    set_minimal_plain_text(shapes["summary"].text_frame, data.summary)

    set_bullet_lines(
        shapes["kompetenzen"].text_frame,
        data.kompetenzen,
        prototypes["kompetenzen"],
    )

    erfahrungen = [(e.label, e.beschreibung) for e in data.relevante_erfahrungen]
    set_categorized_lines(
        shapes["relevante_erfahrungen"].text_frame,
        erfahrungen,
        prototypes["relevante_erfahrungen"],
    )

    set_bullet_lines(
        shapes["ausbildung_karriere"].text_frame,
        data.ausbildung_karriere,
        prototypes["ausbildung_karriere"],
    )
    set_bullet_lines(
        shapes["abschluss_zertifikate"].text_frame,
        data.abschluss_zertifikate,
        prototypes["abschluss_zertifikate"],
    )

    tool_items = [
        ("OSS / Command Management", data.tool_kenntnisse.oss_command_management),
        ("Statistik und Analyse", data.tool_kenntnisse.statistik_analyse),
        ("Planung und Optimierung", data.tool_kenntnisse.planung_optimierung),
        ("Drive Test und Post-Processing", data.tool_kenntnisse.drive_test_post_processing),
        ("Mapping", data.tool_kenntnisse.mapping),
    ]
    tool_items = [(lbl, val) for lbl, val in tool_items if val and val.strip()]
    if tool_items:
        set_categorized_lines(
            shapes["tool_kenntnisse"].text_frame,
            tool_items,
            prototypes["tool_kenntnisse"],
        )
    else:
        clear_shape_text(shapes["tool_kenntnisse"])

    clear_shape_text(shapes["tools_dup"])

    if photo_bytes:
        _replace_photo(slide, photo_bytes, shapes)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path

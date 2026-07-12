"""Extract BeraterprofilData from an existing ORBIT Beraterprofil PPTX."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

from app.schemas.profile import BeraterprofilData, RelevantExperience, ToolKenntnisse
from app.services.content_fit import fit_profile
from app.services.pptx_shapes import TEMPLATE_ANCHORS, find_shape_at

_TOOL_LABEL_MAP = {
    "oss / command management": "oss_command_management",
    "statistik und analyse": "statistik_analyse",
    "planung und optimierung": "planung_optimierung",
    "drive test und post-processing": "drive_test_post_processing",
    "mapping": "mapping",
}


def parse_beraterprofil_pptx(path: str | Path) -> BeraterprofilData:
    prs = Presentation(str(path))
    if not prs.slides:
        raise ValueError("PPTX enthaelt keine Folien")

    slide = prs.slides[0]
    shapes = {key: find_shape_at(slide, left, top) for key, (left, top) in TEMPLATE_ANCHORS.items() if key != "photo" and key != "tools_dup"}

    title = _plain_text(shapes["title"])
    domain = _domain_from_title(title)
    kompetenzen = _bullet_lines(shapes["kompetenzen"])
    erfahrungen = _parse_erfahrungen(shapes["relevante_erfahrungen"])
    ausbildung = _bullet_lines(shapes["ausbildung_karriere"])
    abschluss = _bullet_lines(shapes["abschluss_zertifikate"])
    tools = _parse_tools(shapes["tool_kenntnisse"])

    profile = BeraterprofilData(
        title_domain=domain,
        position=_plain_text(shapes["position"]) or "Consultant",
        schwerpunkte=_plain_text(shapes["schwerpunkte"]),
        summary=_plain_text(shapes["summary"]),
        kompetenzen=kompetenzen or ["Kompetenz 1", "Kompetenz 2", "Kompetenz 3"],
        relevante_erfahrungen=erfahrungen or [
            RelevantExperience(label="Erfahrung", beschreibung="Aus importierter PPTX")
        ],
        ausbildung_karriere=ausbildung or ["Karriere aus importierter PPTX"],
        abschluss_zertifikate=abschluss or ["Abschluss aus importierter PPTX"],
        tool_kenntnisse=tools,
    )
    return fit_profile(profile, preserve_sections=True)


def _plain_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _bullet_lines(shape) -> list[str]:
    if not shape.has_text_frame:
        return []
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    return lines


def _parse_erfahrungen(shape) -> list[RelevantExperience]:
    items: list[RelevantExperience] = []
    for line in _bullet_lines(shape):
        if ": " in line:
            label, desc = line.split(": ", 1)
            items.append(RelevantExperience(label=label.strip(), beschreibung=desc.strip()))
        else:
            items.append(RelevantExperience(label=line, beschreibung=""))
    return items


def _parse_tools(shape) -> ToolKenntnisse:
    tools = ToolKenntnisse()
    for line in _bullet_lines(shape):
        if ": " not in line:
            continue
        category, value = line.split(": ", 1)
        key = _TOOL_LABEL_MAP.get(category.strip().lower())
        if key:
            setattr(tools, key, value.strip())
    return tools


def _domain_from_title(title: str) -> str:
    for sep in ("–", "-", "—"):
        if sep in title:
            part = title.split(sep, 1)[1].strip()
            if part:
                return part
    return "IT-Beratung"
